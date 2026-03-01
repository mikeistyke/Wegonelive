from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT_PATH = r"c:\Users\Mike\OneDrive\Mikes Data\Desktop\WeGoneLive\My Docs\client_safe_architecture_overview.pptx"

BG = RGBColor(11, 18, 32)
FG = RGBColor(235, 241, 255)
MUTED = RGBColor(160, 177, 211)
ACCENT = RGBColor(39, 203, 166)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def style_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title, subtitle=None):
    t = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(1.2))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Segoe UI Semibold"
    p.font.size = Pt(34)
    p.font.color.rgb = FG

    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.7))
        stf = s.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Segoe UI"
        sp.font.size = Pt(16)
        sp.font.color.rgb = MUTED


def add_bullets(slide, lines, x=0.95, y=2.0, w=11.8, h=4.9, size=23):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = FG
        p.space_after = Pt(16)


def add_two_column(slide, left_title, left_lines, right_title, right_lines):
    # left panel
    lt = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.9), Inches(0.6))
    ltf = lt.text_frame
    ltf.text = left_title
    ltf.paragraphs[0].font.name = "Segoe UI Semibold"
    ltf.paragraphs[0].font.size = Pt(20)
    ltf.paragraphs[0].font.color.rgb = ACCENT

    lb = slide.shapes.add_textbox(Inches(0.9), Inches(2.45), Inches(5.8), Inches(3.9))
    lbtf = lb.text_frame
    lbtf.clear()
    for i, line in enumerate(left_lines):
        p = lbtf.paragraphs[0] if i == 0 else lbtf.add_paragraph()
        p.text = line
        p.font.name = "Segoe UI"
        p.font.size = Pt(18)
        p.font.color.rgb = FG
        p.space_after = Pt(10)

    # right panel
    rt = slide.shapes.add_textbox(Inches(7.0), Inches(1.9), Inches(5.4), Inches(0.6))
    rtf = rt.text_frame
    rtf.text = right_title
    rtf.paragraphs[0].font.name = "Segoe UI Semibold"
    rtf.paragraphs[0].font.size = Pt(20)
    rtf.paragraphs[0].font.color.rgb = ACCENT

    rb = slide.shapes.add_textbox(Inches(7.1), Inches(2.45), Inches(5.2), Inches(3.9))
    rbtf = rb.text_frame
    rbtf.clear()
    for i, line in enumerate(right_lines):
        p = rbtf.paragraphs[0] if i == 0 else rbtf.add_paragraph()
        p.text = line
        p.font.name = "Segoe UI"
        p.font.size = Pt(18)
        p.font.color.rgb = FG
        p.space_after = Pt(10)


# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
style_background(slide)
add_title(slide, "Live Commerce Architecture", "Client-safe overview · Designed for confidence without exposing proprietary mechanics")
add_bullets(slide, [
    "Operationally strong: real-time event controls + fallback decision paths",
    "Governed by clear policy modules and auditable event records",
    "Built for repeatable execution across live sessions"
], y=2.25, size=24)

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
style_background(slide)
add_title(slide, "What Clients Receive", "Business value, not internal secret sauce")
add_bullets(slide, [
    "Reliable live session orchestration",
    "Pricing guardrails and fallback handling",
    "Role-based operational dashboard",
    "Post-event traceability and performance reporting"
], y=2.0, size=22)

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
style_background(slide)
add_title(slide, "System Layers")
add_two_column(
    slide,
    "Execution Plane",
    [
        "• Presentation layer (public + operator views)",
        "• Real-time bidding and event control services",
        "• Decision orchestration modules"
    ],
    "Intelligence & Governance",
    [
        "• Valuation integration adapters",
        "• Metrics and analytics pipeline",
        "• Identity, access, and audit trail"
    ]
)

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
style_background(slide)
add_title(slide, "Runtime Flow")
add_bullets(slide, [
    "1) Session configured with inventory and valuation references",
    "2) Live event executes with controlled bid intake",
    "3) Policy evaluation at close selects outcome path",
    "4) Checkout or fallback offer orchestration",
    "5) Metrics + notices written for post-run review"
], y=2.0, size=20)

# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
style_background(slide)
add_title(slide, "Control Model")
add_two_column(
    slide,
    "Operator Controls",
    [
        "• Session activation",
        "• Item/row targeting",
        "• Guided fallback actions",
        "• Single-pane monitoring"
    ],
    "Risk Controls",
    [
        "• Policy thresholds",
        "• Soft/hard validation options",
        "• Restricted internal insights",
        "• Event-level traceability"
    ]
)

# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
style_background(slide)
add_title(slide, "Architecture Gems")
add_bullets(slide, [
    "Policy-first orchestration keeps outcomes consistent under pressure",
    "Fallback path protects value when primary close conditions are not met",
    "Row-level targeting keeps live execution aligned with inventory state",
    "Audit-friendly notices and metrics reduce post-event ambiguity"
], y=2.0, size=20)

# Slide 7
slide = prs.slides.add_slide(prs.slide_layouts[6])
style_background(slide)
add_title(slide, "Security & Compliance Posture")
add_bullets(slide, [
    "Least-privilege access for internal operational surfaces",
    "Secrets and privileged keys remain server-side",
    "Separation of public experience and restricted controls",
    "Traceable decisions for governance and replay"
], y=2.0, size=21)

# Slide 8
slide = prs.slides.add_slide(prs.slide_layouts[6])
style_background(slide)
add_title(slide, "Scalability & Reliability")
add_bullets(slide, [
    "Modular architecture supports phased expansion",
    "Near real-time visibility for operators during event execution",
    "Graceful fallback behavior if optional integrations are unavailable",
    "Deployment strategy can align with client cloud standards"
], y=2.0, size=21)

# Slide 9
slide = prs.slides.add_slide(prs.slide_layouts[6])
style_background(slide)
add_title(slide, "Engagement Model")
add_bullets(slide, [
    "Phase 1: Discovery + policy mapping",
    "Phase 2: Pilot with measured success criteria",
    "Phase 3: Controlled rollout with operator enablement",
    "NDA deep-dive available for implementation specifics"
], y=2.0, size=21)

# Slide 10
slide = prs.slides.add_slide(prs.slide_layouts[6])
style_background(slide)
add_title(slide, "Closing")
add_bullets(slide, [
    "This architecture is built to execute live commerce with discipline.",
    "It combines operational clarity, policy guardrails, and measurable outcomes.",
    "Next step: align pilot scope to your event model and governance requirements."
], y=2.35, size=22)

prs.save(OUTPUT_PATH)
print(OUTPUT_PATH)
