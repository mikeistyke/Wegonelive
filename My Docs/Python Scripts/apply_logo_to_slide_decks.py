from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

DECK_DIR = Path(r"c:\Users\Mike\OneDrive\Mikes Data\Desktop\WeGoneLive\My Docs\Slide Deck")
LOGO_PATH = Path(r"c:\Users\Mike\OneDrive\Mikes Data\Desktop\WeGoneLive\My Docs\Logo\wegonelive_logo_only.jpg")
LOGO_SHAPE_NAME = "WGL_Logo"


def apply_logo(deck_path: Path) -> tuple[int, int]:
    prs = Presentation(str(deck_path))
    added = 0
    skipped = 0

    for slide in prs.slides:
        has_logo = any(getattr(shape, "name", "") == LOGO_SHAPE_NAME for shape in slide.shapes)
        if has_logo:
            skipped += 1
            continue

        logo_width = Inches(1.45)
        left = prs.slide_width - logo_width - Inches(0.30)
        top = Inches(0.20)

        logo_shape = slide.shapes.add_picture(str(LOGO_PATH), left, top, width=logo_width)
        logo_shape.name = LOGO_SHAPE_NAME
        added += 1

    prs.save(str(deck_path))
    return added, skipped


def main() -> None:
    if not DECK_DIR.exists():
        raise FileNotFoundError(f"Deck folder not found: {DECK_DIR}")
    if not LOGO_PATH.exists():
        raise FileNotFoundError(f"Logo image not found: {LOGO_PATH}")

    decks = sorted(DECK_DIR.glob("*.pptx"))
    if not decks:
        print("No PowerPoint files found.")
        return

    for deck in decks:
        added, skipped = apply_logo(deck)
        print(f"{deck.name}: added={added}, already_present={skipped}")


if __name__ == "__main__":
    main()
