from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT_PATH = r"c:\Users\Mike\OneDrive\Mikes Data\Desktop\WeGoneLive\My Docs\client_vertical_collectibles_fashion_showcase.pptx"

BG = RGBColor(10, 14, 28)
PANEL = RGBColor(18, 26, 44)
FG = RGBColor(240, 246, 255)
MUTED = RGBColor(164, 182, 216)
ACCENT = RGBColor(66, 233, 189)
GOLD = RGBColor(245, 196, 93)
ROSE = RGBColor(244, 126, 162)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def style_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_header(slide, title, subtitle=""):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12.0), Inches(1.1))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Segoe UI Semibold"
    p.font.size = Pt(35)
    p.font.color.rgb = FG

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.72), Inches(1.18), Inches(11.9), Inches(0.55))
        stf = sub.text_frame
        stf.text = subtitle
        stf.paragraphs[0].font.name = "Segoe UI"
        stf.paragraphs[0].font.size = Pt(15)
        stf.paragraphs[0].font.color.rgb = MUTED


def add_banner(slide, text):
    shape = slide.shapes.add_shape(1, Inches(0.7), Inches(1.72), Inches(12.0), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = ACCENT
    tf = shape.text_frame
    tf.text = text
    tf.paragraphs[0].font.name = "Segoe UI"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = ACCENT


def add_bullets(slide, lines, y=2.3, size=21):
    box = slide.shapes.add_textbox(Inches(0.85), Inches(y), Inches(11.8), Inches(4.8))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = FG
        p.space_after = Pt(14)


def add_dual_vertical_panel(slide, left_title, left_color, left_points, right_title, right_color, right_points):
    left = slide.shapes.add_shape(1, Inches(0.75), Inches(2.0), Inches(5.95), Inches(4.9))
    left.fill.solid()
    left.fill.fore_color.rgb = PANEL
    left.line.color.rgb = RGBColor(37, 55, 87)

    lt = slide.shapes.add_textbox(Inches(1.02), Inches(2.22), Inches(5.3), Inches(0.52))
    ltf = lt.text_frame
    ltf.text = left_title
    ltf.paragraphs[0].font.name = "Segoe UI Semibold"
    ltf.paragraphs[0].font.size = Pt(21)
    ltf.paragraphs[0].font.color.rgb = left_color

    lb = slide.shapes.add_textbox(Inches(1.04), Inches(2.8), Inches(5.2), Inches(3.8))
    lbtf = lb.text_frame
    lbtf.clear()
    for i, line in enumerate(left_points):
        p = lbtf.paragraphs[0] if i == 0 else lbtf.add_paragraph()
        p.text = line
        p.font.name = "Segoe UI"
        p.font.size = Pt(16)
        p.font.color.rgb = FG
        p.space_after = Pt(9)

    right = slide.shapes.add_shape(1, Inches(6.95), Inches(2.0), Inches(5.65), Inches(4.9))
    right.fill.solid()
    right.fill.fore_color.rgb = PANEL
    right.line.color.rgb = RGBColor(37, 55, 87)

    rt = slide.shapes.add_textbox(Inches(7.22), Inches(2.22), Inches(5.1), Inches(0.52))
    rtf = rt.text_frame
    rtf.text = right_title
    rtf.paragraphs[0].font.name = "Segoe UI Semibold"
    rtf.paragraphs[0].font.size = Pt(21)
    rtf.paragraphs[0].font.color.rgb = right_color

    rb = slide.shapes.add_textbox(Inches(7.24), Inches(2.8), Inches(5.0), Inches(3.8))
    rbtf = rb.text_frame
    rbtf.clear()
    for i, line in enumerate(right_points):
        p = rbtf.paragraphs[0] if i == 0 else rbtf.add_paragraph()
        p.text = line
        p.font.name = "Segoe UI"
        p.font.size = Pt(16)
        p.font.color.rgb = FG
        p.space_after = Pt(9)


# Slide 1
s = prs.slides.add_slide(prs.slide_layouts[6])
style_bg(s)
add_header(s, "Client Vertical Showcase", "Collectibles (Coins & Stamps) + Independent Fashion & Design")
add_banner(s, "Professional by design: elevated client narrative with proprietary logic intentionally abstracted")
add_bullets(s, [
    "Built for high-variance inventory where trust, speed, and control must coexist",
    "Supports premium storytelling without sacrificing operational discipline",
    "Turns complex close decisions into clear, auditable operator actions",
], y=2.42, size=22)

# Slide 2
s = prs.slides.add_slide(prs.slide_layouts[6])
style_bg(s)
add_header(s, "Vertical Fit")
add_dual_vertical_panel(
    s,
    "Collectibles: Coins & Stamps",
    GOLD,
    [
        "• Condition and rarity-sensitive pricing moments",
        "• Value-protection logic for uncertain close outcomes",
        "• Traceable decisions for buyer confidence"
    ],
    "Independent Fashion & Design",
    ROSE,
    [
        "• One-off drops and non-store inventory",
        "• Live momentum with controlled fallback paths",
        "• Creator-friendly operator workflows"
    ]
)

# Slide 3
s = prs.slides.add_slide(prs.slide_layouts[6])
style_bg(s)
add_header(s, "Experience Architecture")
add_bullets(s, [
    "1) Curated session setup with inventory + valuation references",
    "2) Real-time event execution with operator control surfaces",
    "3) Policy-guided close path selection",
    "4) Structured fallback offer sequence when required",
    "5) Metrics and notice capture for post-event replay"
], y=2.1, size=20)

# Slide 4
s = prs.slides.add_slide(prs.slide_layouts[6])
style_bg(s)
add_header(s, "What Prospects Can See")
add_dual_vertical_panel(
    s,
    "Visible Strengths",
    ACCENT,
    [
        "• Session-level row targeting",
        "• Clear in-event decision actions",
        "• Internal metrics for operators",
        "• Reliable event traceability"
    ],
    "Protected Internals",
    GOLD,
    [
        "• Threshold tuning mechanics",
        "• Sequencing heuristics",
        "• Optimization strategy",
        "• Detailed internals under NDA"
    ]
)

# Slide 5
s = prs.slides.add_slide(prs.slide_layouts[6])
style_bg(s)
add_header(s, "Business Outcomes")
add_bullets(s, [
    "Higher operator confidence during complex live close moments",
    "Reduced margin leakage through policy-governed fallback decisions",
    "Stronger buyer trust from consistent and transparent event handling",
    "Cleaner post-event coaching via structured event records"
], y=2.2, size=22)

# Slide 6
s = prs.slides.add_slide(prs.slide_layouts[6])
style_bg(s)
add_header(s, "Pilot Blueprint")
add_dual_vertical_panel(
    s,
    "Pilot Scope",
    ACCENT,
    [
        "• 1 curated session type",
        "• 1 operator team",
        "• 2 scenario paths (direct + fallback)",
        "• Defined success scorecard"
    ],
    "Readiness Gate",
    GOLD,
    [
        "• Go/No-Go runbook completion",
        "• Decision quality checks",
        "• Trace + metrics review",
        "• Rollout recommendation"
    ]
)

# Slide 7
s = prs.slides.add_slide(prs.slide_layouts[6])
style_bg(s)
add_header(s, "Close")
add_bullets(s, [
    "This platform supports premium live selling in categories where every close decision matters.",
    "It delivers controlled execution, elegant operator flow, and confidence-grade governance.",
    "Next step: align pilot design to your inventory profile and event cadence."
], y=2.45, size=22)

prs.save(OUTPUT_PATH)
print(OUTPUT_PATH)
