from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT_PATH = r"c:\Users\Mike\OneDrive\Mikes Data\Desktop\WeGoneLive\My Docs\client_vertical_luxury_resale_showcase.pptx"

BG = RGBColor(7, 12, 24)
PANEL = RGBColor(16, 24, 42)
FG = RGBColor(241, 245, 255)
MUTED = RGBColor(163, 180, 211)
ACCENT = RGBColor(56, 226, 179)
GOLD = RGBColor(244, 195, 87)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def style_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_header(slide, title, subtitle=""):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12.0), Inches(1.1))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Segoe UI Semibold"
    p.font.size = Pt(36)
    p.font.color.rgb = FG

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.72), Inches(1.18), Inches(11.9), Inches(0.55))
        stf = sub.text_frame
        stf.text = subtitle
        stf.paragraphs[0].font.name = "Segoe UI"
        stf.paragraphs[0].font.size = Pt(15)
        stf.paragraphs[0].font.color.rgb = MUTED


def add_banner(slide, text):
    shape = slide.shapes.add_shape(1, Inches(0.7), Inches(1.72), Inches(12.0), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = ACCENT
    tf = shape.text_frame
    tf.text = text
    tf.paragraphs[0].font.name = "Segoe UI"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = ACCENT


def add_bullets(slide, lines, y=2.35, size=22):
    box = slide.shapes.add_textbox(Inches(0.85), Inches(y), Inches(11.8), Inches(4.7))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = FG
        p.space_after = Pt(14)


def add_two_panels(slide, left_title, left_lines, right_title, right_lines):
    left = slide.shapes.add_shape(1, Inches(0.75), Inches(2.0), Inches(5.95), Inches(4.8))
    left.fill.solid()
    left.fill.fore_color.rgb = PANEL
    left.line.color.rgb = RGBColor(36, 53, 84)

    ltitle = slide.shapes.add_textbox(Inches(1.05), Inches(2.25), Inches(5.3), Inches(0.45))
    ltitle.text_frame.text = left_title
    ltitle.text_frame.paragraphs[0].font.name = "Segoe UI Semibold"
    ltitle.text_frame.paragraphs[0].font.size = Pt(20)
    ltitle.text_frame.paragraphs[0].font.color.rgb = ACCENT

    lbody = slide.shapes.add_textbox(Inches(1.05), Inches(2.75), Inches(5.2), Inches(3.8))
    ltf = lbody.text_frame
    ltf.clear()
    for i, line in enumerate(left_lines):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = line
        p.font.name = "Segoe UI"
        p.font.size = Pt(16)
        p.font.color.rgb = FG
        p.space_after = Pt(10)

    right = slide.shapes.add_shape(1, Inches(6.95), Inches(2.0), Inches(5.65), Inches(4.8))
    right.fill.solid()
    right.fill.fore_color.rgb = PANEL
    right.line.color.rgb = RGBColor(36, 53, 84)

    rtitle = slide.shapes.add_textbox(Inches(7.25), Inches(2.25), Inches(5.1), Inches(0.45))
    rtitle.text_frame.text = right_title
    rtitle.text_frame.paragraphs[0].font.name = "Segoe UI Semibold"
    rtitle.text_frame.paragraphs[0].font.size = Pt(20)
    rtitle.text_frame.paragraphs[0].font.color.rgb = GOLD

    rbody = slide.shapes.add_textbox(Inches(7.25), Inches(2.75), Inches(5.0), Inches(3.8))
    rtf = rbody.text_frame
    rtf.clear()
    for i, line in enumerate(right_lines):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = line
        p.font.name = "Segoe UI"
        p.font.size = Pt(16)
        p.font.color.rgb = FG
        p.space_after = Pt(10)


# Slide 1
s = prs.slides.add_slide(prs.slide_layouts[6])
style_bg(s)
add_header(s, "Luxury Resale Live Commerce", "Showcase Edition · Professional narrative for client-facing demos")
add_banner(s, "Positioning: premium experience orchestration with disciplined control and measurable outcomes")
add_bullets(s, [
    "Built for high-attention products where confidence and control matter",
    "Balances conversion speed with value-protection decision paths",
    "Delivers operator clarity during live moments and clean post-event traceability",
], y=2.45, size=23)

# Slide 2
s = prs.slides.add_slide(prs.slide_layouts[6])
style_bg(s)
add_header(s, "Why This Matters for Premium Catalogs")
add_two_panels(
    s,
    "Commercial Pressure",
    [
        "• Live events move fast under uncertain pricing",
        "• Operators need certainty without slowing momentum",
        "• Missed decisions create margin leakage"
    ],
    "Platform Response",
    [
        "• Policy-driven close decisions",
        "• Structured fallback orchestration",
        "• Traceable outcomes for every critical action"
    ]
)

# Slide 3
s = prs.slides.add_slide(prs.slide_layouts[6])
style_bg(s)
add_header(s, "Architecture in Motion")
add_bullets(s, [
    "1) Session + inventory context loaded before event start",
    "2) Live bid flow runs with operator controls and guardrails",
    "3) Close decision module selects checkout or fallback path",
    "4) Outcome syncs into internal metrics + notice trail",
    "5) Post-event review is replayable and auditable"
], y=2.2, size=20)

# Slide 4
s = prs.slides.add_slide(prs.slide_layouts[6])
style_bg(s)
add_header(s, "Client-Safe Feature Highlights")
add_two_panels(
    s,
    "Execution Highlights",
    [
        "• Session-scoped item targeting",
        "• Real-time operator signal layer",
        "• Guided decision actions in-event",
        "• Restricted internal insight surfaces"
    ],
    "Value Highlights",
    [
        "• Better decision consistency",
        "• Reduced revenue leakage",
        "• Faster operator response under load",
        "• Cleaner review and coaching loop"
    ]
)

# Slide 5
s = prs.slides.add_slide(prs.slide_layouts[6])
style_bg(s)
add_header(s, "Governance Without Friction")
add_bullets(s, [
    "Least-privilege controls for internal operational routes",
    "Server-side protection for sensitive credentials and policy execution",
    "Event-level trace capture for accountable decisions",
    "Separation of public experience and private control surfaces"
], y=2.2, size=21)

# Slide 6
s = prs.slides.add_slide(prs.slide_layouts[6])
style_bg(s)
add_header(s, "What We Keep Private (By Design)")
add_bullets(s, [
    "Internal threshold calibration and tuning mechanics",
    "Fallback sequencing heuristics and optimization patterns",
    "Prompting/orchestration internals and model strategies",
    "Implementation details available under NDA only"
], y=2.2, size=21)

# Slide 7
s = prs.slides.add_slide(prs.slide_layouts[6])
style_bg(s)
add_header(s, "Pilot Path")
add_two_panels(
    s,
    "Phase 1 — Alignment",
    [
        "• event model mapping",
        "• policy constraints",
        "• success criteria"
    ],
    "Phase 2 — Controlled Pilot",
    [
        "• selected catalog segment",
        "• operator enablement",
        "• go/no-go review"
    ]
)

# Slide 8
s = prs.slides.add_slide(prs.slide_layouts[6])
style_bg(s)
add_header(s, "Close")
add_bullets(s, [
    "This architecture is designed to feel premium in front of buyers and disciplined behind the scenes.",
    "It scales from pilot to production while protecting operational confidence.",
    "Next step: tailor pilot scope to your category and event cadence."
], y=2.4, size=22)

prs.save(OUTPUT_PATH)
print(OUTPUT_PATH)
